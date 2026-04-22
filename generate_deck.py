# pyright: reportMissingImports=false, reportMissingModuleSource=false
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

VENDOR_DIR = Path(__file__).resolve().parent / ".vendor"
if VENDOR_DIR.exists():
    sys.path.insert(0, str(VENDOR_DIR))

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


COLOR_ALIASES = {
    "title": (22, 33, 62),
    "accent": (43, 108, 176),
    "secondary": (79, 93, 117),
    "light_fill": (236, 242, 248),
    "pale_green": (230, 245, 234),
    "pale_orange": (253, 237, 222),
    "white": (255, 255, 255),
    "black": (0, 0, 0),
}

ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

SHAPES = {
    "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
}


def rgb_color(value) -> RGBColor:
    if isinstance(value, RGBColor):
        return value
    if value is None:
        value = "white"
    if isinstance(value, str):
        alias = COLOR_ALIASES.get(value.lower())
        if alias:
            return RGBColor(*alias)
        hex_value = value.lstrip("#")
        if len(hex_value) == 6 and re.fullmatch(r"[0-9A-Fa-f]{6}", hex_value):
            return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))
    if isinstance(value, list) and len(value) == 3:
        return RGBColor(*[int(part) for part in value])
    raise ValueError(f"Unsupported color value: {value!r}")


def inches(value) -> int:
    return Inches(float(value))


def points(value) -> Pt:
    return Pt(float(value))


def set_fill(fill, color) -> None:
    fill.solid()
    fill.fore_color.rgb = rgb_color(color)


def set_background(slide, color="white") -> None:
    set_fill(slide.background.fill, color)


def set_font(run, spec: dict, defaults: dict | None = None) -> None:
    defaults = defaults or {}
    run.font.name = spec.get("font_name", defaults.get("font_name", "Aptos"))
    run.font.size = points(spec.get("font_size", defaults.get("font_size", 18)))
    run.font.bold = spec.get("bold", defaults.get("bold", False))
    run.font.italic = spec.get("italic", defaults.get("italic", False))
    run.font.color.rgb = rgb_color(spec.get("color", defaults.get("color", "black")))


def add_textbox(slide, spec: dict, defaults: dict | None = None) -> None:
    defaults = defaults or {}
    shape = slide.shapes.add_textbox(
        inches(spec.get("x", defaults.get("x", 0.6))),
        inches(spec.get("y", defaults.get("y", 0.6))),
        inches(spec.get("w", defaults.get("w", 12.0))),
        inches(spec.get("h", defaults.get("h", 0.8))),
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = spec.get("word_wrap", defaults.get("word_wrap", True))

    paragraphs = spec.get("paragraphs")
    if paragraphs is None:
        if "bullets" in spec:
            paragraphs = [{"text": bullet, "level": 0} for bullet in spec["bullets"]]
        else:
            paragraphs = [{"text": spec.get("text", ""), "level": spec.get("level", 0)}]

    for idx, paragraph_spec in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.level = int(paragraph_spec.get("level", 0))
        alignment_name = paragraph_spec.get("alignment", spec.get("alignment", defaults.get("alignment", "left")))
        paragraph.alignment = ALIGNMENTS[alignment_name]
        paragraph.space_after = points(paragraph_spec.get("space_after", spec.get("space_after", defaults.get("space_after", 0))))
        run = paragraph.add_run()
        run.text = paragraph_spec["text"]
        set_font(run, paragraph_spec, {**defaults, **spec})


def add_box(slide, spec: dict) -> None:
    shape = slide.shapes.add_shape(
        SHAPES[spec.get("shape", "rounded_rectangle")],
        inches(spec["x"]),
        inches(spec["y"]),
        inches(spec["w"]),
        inches(spec["h"]),
    )
    set_fill(shape.fill, spec.get("fill", "light_fill"))
    shape.line.color.rgb = rgb_color(spec.get("line_color", "accent"))
    shape.line.width = points(spec.get("line_width", 1.25))

    tf = shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.alignment = ALIGNMENTS[spec.get("alignment", "center")]
    run = paragraph.add_run()
    run.text = spec["text"]
    set_font(
        run,
        spec,
        {
            "font_name": "Aptos",
            "font_size": 15,
            "bold": False,
            "color": spec.get("text_color", "title"),
        },
    )


def add_arrow(slide, spec: dict) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        inches(spec["x1"]),
        inches(spec["y1"]),
        inches(spec["x2"]),
        inches(spec["y2"]),
    )
    line.line.color.rgb = rgb_color(spec.get("color", "secondary"))
    line.line.width = points(spec.get("width", 2))
    if spec.get("end_arrow", True):
        line.line.end_arrowhead = True


def add_banner(slide, spec: dict, prs: Presentation) -> None:
    width = spec.get("w", prs.slide_width / Inches(1))
    shape = slide.shapes.add_shape(
        SHAPES[spec.get("shape", "rectangle")],
        inches(spec.get("x", 0)),
        inches(spec.get("y", 0)),
        inches(width),
        inches(spec["h"]),
    )
    set_fill(shape.fill, spec.get("fill", "title"))
    if spec.get("hide_line", True):
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb_color(spec.get("line_color", "title"))


def style_table_cell(cell, *, fill_color, text_color, font_size, bold=False, font_name="Aptos") -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb_color(fill_color)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = points(font_size)
            run.font.bold = bold
            run.font.color.rgb = rgb_color(text_color)


def add_table(slide, spec: dict) -> None:
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    total_rows = len(rows) + (1 if headers else 0)
    total_cols = len(headers) if headers else max(len(row) for row in rows)
    table = slide.shapes.add_table(
        total_rows,
        total_cols,
        inches(spec["x"]),
        inches(spec["y"]),
        inches(spec["w"]),
        inches(spec["h"]),
    ).table

    for index, width in enumerate(spec.get("column_widths", [])):
        table.columns[index].width = inches(width)

    row_offset = 0
    if headers:
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            style_table_cell(
                cell,
                fill_color=spec.get("header_fill", "title"),
                text_color=spec.get("header_text_color", "white"),
                font_size=spec.get("header_font_size", 14),
                bold=True,
            )
        row_offset = 1

    alternating_fills = spec.get("row_fills", ["white", "light_fill"])
    for row_idx, row in enumerate(rows):
        fill_color = alternating_fills[row_idx % len(alternating_fills)]
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + row_offset, col_idx)
            cell.text = str(value)
            style_table_cell(
                cell,
                fill_color=fill_color,
                text_color=spec.get("text_color", "black"),
                font_size=spec.get("font_size", 13),
            )


def render_element(slide, spec: dict, prs: Presentation) -> None:
    kind = spec["kind"]
    if kind == "title":
        add_textbox(
            slide,
            spec,
            {"x": 0.6, "y": 0.35, "w": 12.0, "h": 0.8, "font_name": "Aptos Display", "font_size": 28, "bold": True, "color": "title"},
        )
    elif kind == "subtitle":
        add_textbox(
            slide,
            spec,
            {"x": 0.6, "y": 0.95, "w": 12.0, "h": 0.45, "font_name": "Aptos", "font_size": 12, "color": "secondary"},
        )
    elif kind == "textbox":
        add_textbox(slide, spec)
    elif kind == "bullets":
        add_textbox(
            slide,
            spec,
            {"font_name": "Aptos", "font_size": 20, "color": "black", "space_after": 10},
        )
    elif kind == "note":
        add_textbox(
            slide,
            spec,
            {"font_name": "Aptos", "font_size": 11, "italic": True, "color": "secondary", "alignment": "center"},
        )
    elif kind == "box":
        add_box(slide, spec)
    elif kind == "arrow":
        add_arrow(slide, spec)
    elif kind == "banner":
        add_banner(slide, spec, prs)
    elif kind == "table":
        add_table(slide, spec)
    else:
        raise ValueError(f"Unsupported element kind: {kind}")


def slide_sort_key(path: Path) -> tuple[int, str]:
    match = re.search(r"(\d+)", path.stem)
    return (int(match.group(1)) if match else 10**9, path.name)


def load_slide_specs(slides_dir: Path) -> list[dict]:
    slide_paths = sorted(slides_dir.glob("slide*.yaml"), key=slide_sort_key)
    if not slide_paths:
        raise SystemExit(f"No slide YAML files found in {slides_dir}")
    return [yaml.safe_load(path.read_text()) for path in slide_paths]


def build_presentation(slides_dir: Path) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_spec in load_slide_specs(slides_dir):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, slide_spec.get("background", "white"))
        for element in slide_spec.get("elements", []):
            render_element(slide, element, prs)

    return prs


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a PowerPoint deck from slide YAML files.")
    parser.add_argument("--slides-dir", default="slides", help="Directory containing slide1.yaml, slide2.yaml, etc.")
    parser.add_argument("--output", default="opentelemetry-signoz-observability-stack.pptx", help="Output .pptx path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    slides_dir = Path(args.slides_dir)
    output = Path(args.output)
    prs = build_presentation(slides_dir)
    prs.save(output)
    print(f"Wrote {output.resolve()} from {slides_dir.resolve()}")


if __name__ == "__main__":
    main()
